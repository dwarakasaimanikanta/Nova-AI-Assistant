"""
utils/document_manager.py
-------------------------
DocumentManager manages Retrieval-Augmented Generation (RAG) capabilities,
supporting text extractions from PDF, DOCX, XLSX, PPTX, TXT, and Markdown files,
chunking, embedding generation (Gemini/Ollama), and semantic similarity search in SQLite.
"""

import os
import json
import sqlite3
from pathlib import Path
from typing import Any, List, Dict, Tuple, Optional
import requests
from config import DATA_DIR, OLLAMA_HOST, GEMINI_API_KEY
from utils.logger import get_logger

logger = get_logger(__name__)


class DocumentManager:
    """Manages document text extraction, chunk parsing, vector database index updates, and searches."""

    def __init__(self, db_path: Optional[Path] = None) -> None:
        """
        Initialize the DocumentManager.

        Args:
            db_path: Path to the SQLite database. Defaults to data/long_term_memory.db.
        """
        self.db_path = db_path or (DATA_DIR / "long_term_memory.db")
        self._init_db()

    def _init_db(self) -> None:
        """Initializes the document_chunks schema inside the database."""
        self.db_path.parent.mkdir(parents=True, exist_ok=True)
        try:
            conn = sqlite3.connect(str(self.db_path), check_same_thread=False)
            cursor = conn.cursor()
            cursor.execute("""
                CREATE TABLE IF NOT EXISTS document_chunks (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    doc_name TEXT,
                    chunk_index INTEGER,
                    content TEXT,
                    embedding TEXT
                )
            """)
            conn.commit()
            conn.close()
            logger.info("RAG document_chunks database initialized successfully.")
        except Exception as e:
            logger.error("Failed to initialize document table in SQLite: %s", e)

    # --- Text Extraction Parsers ---

    def extract_text(self, filepath: Path) -> str:
        """
        Extracts raw text content from the target document path.

        Args:
            filepath: Path to the document.

        Returns:
            The raw text string extracted from the file.
        """
        if not filepath.exists():
            raise FileNotFoundError(f"File not found: {filepath}")

        suffix = filepath.suffix.lower()

        if suffix == ".pdf":
            return self._extract_pdf(filepath)
        elif suffix == ".docx":
            return self._extract_docx(filepath)
        elif suffix == ".xlsx":
            return self._extract_xlsx(filepath)
        elif suffix == ".pptx":
            return self._extract_pptx(filepath)
        elif suffix in (".txt", ".md"):
            return self._extract_txt_md(filepath)
        else:
            raise ValueError(f"Unsupported document file extension: '{suffix}'")

    def _extract_pdf(self, filepath: Path) -> str:
        import pypdf
        try:
            reader = pypdf.PdfReader(filepath)
            text_lines = []
            for idx, page in enumerate(reader.pages):
                extracted = page.extract_text()
                if extracted:
                    text_lines.append(extracted)
            return "\n".join(text_lines)
        except Exception as e:
            logger.error("PDF parsing error on '%s': %s", filepath.name, e)
            raise RuntimeError(f"PDF extraction error: {e}")

    def _extract_docx(self, filepath: Path) -> str:
        import docx
        try:
            doc = docx.Document(filepath)
            paragraphs = [p.text for p in doc.paragraphs if p.text]
            return "\n".join(paragraphs)
        except Exception as e:
            logger.error("Word document parsing error on '%s': %s", filepath.name, e)
            raise RuntimeError(f"Word document extraction error: {e}")

    def _extract_xlsx(self, filepath: Path) -> str:
        import openpyxl
        try:
            wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
            lines = []
            for sheetname in wb.sheetnames:
                sheet = wb[sheetname]
                lines.append(f"--- Sheet: {sheetname} ---")
                for row in sheet.iter_rows(values_only=True):
                    if any(row):
                        row_str = ", ".join([str(cell) for cell in row if cell is not None])
                        lines.append(row_str)
            return "\n".join(lines)
        except Exception as e:
            logger.error("Excel document parsing error on '%s': %s", filepath.name, e)
            raise RuntimeError(f"Excel extraction error: {e}")

    def _extract_pptx(self, filepath: Path) -> str:
        from pptx import Presentation
        try:
            prs = Presentation(filepath)
            lines = []
            for i, slide in enumerate(prs.slides):
                lines.append(f"--- Slide {i+1} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        lines.append(shape.text.strip())
            return "\n".join(lines)
        except Exception as e:
            logger.error("Powerpoint document parsing error on '%s': %s", filepath.name, e)
            raise RuntimeError(f"PowerPoint extraction error: {e}")

    def _extract_txt_md(self, filepath: Path) -> str:
        try:
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        except Exception as e:
            logger.error("Text file read error on '%s': %s", filepath.name, e)
            raise RuntimeError(f"Text read error: {e}")

    # --- Chunking & Embedding Helpers ---

    def chunk_text(self, text: str, chunk_size: int = 500, overlap: int = 100) -> List[str]:
        """Splits raw text into sliding window chunks."""
        chunks = []
        if not text:
            return chunks

        start = 0
        text_len = len(text)
        while start < text_len:
            end = min(start + chunk_size, text_len)
            chunks.append(text[start:end].strip())
            if end >= text_len:
                break
            start += chunk_size - overlap
        return [c for c in chunks if c]

    def _get_embedding(self, text: str) -> List[float]:
        """Fetches vector embedding. Uses Gemini online, Ollama offline, or returns dummy list."""
        # 1. Try Gemini
        if GEMINI_API_KEY and os.getenv("ENVIRONMENT") != "test":
            try:
                import google.generativeai as genai
                genai.configure(api_key=GEMINI_API_KEY)
                res = genai.embed_content(
                    model="models/gemini-embedding-001",
                    content=text,
                    task_type="retrieval_document"
                )
                embedding = res.get("embedding")
                if embedding:
                    return embedding
            except Exception as e:
                logger.warning("Failed to retrieve Gemini embedding: %s. Trying Ollama fallback.", e)

        # 2. Try Ollama Offline
        if os.getenv("ENVIRONMENT") != "test":
            try:
                url = f"{OLLAMA_HOST}/api/embeddings"
                response = requests.post(url, json={"model": "llama3", "prompt": text}, timeout=3.0)
                if response.status_code == 200:
                    embedding = response.json().get("embedding")
                    if embedding:
                        return embedding
            except Exception as e:
                logger.debug("Failed to retrieve local Ollama embedding: %s", e)

        # 3. Headless Mock Fallback (Fixed dimensions vector)
        return [0.01 * (i % 10) for i in range(768)]

    # --- SQLite CRUD Operations ---

    def ingest_document(self, filepath: Path) -> int:
        """Parses document, chunks text, generates embeddings, and saves to database."""
        text = self.extract_text(filepath)
        chunks = self.chunk_text(text)
        doc_name = filepath.name

        conn = sqlite3.connect(str(self.db_path), check_same_thread=False)
        cursor = conn.cursor()

        # Delete any existing document index records to avoid duplication
        cursor.execute("DELETE FROM document_chunks WHERE doc_name = ?", (doc_name,))

        count = 0
        for idx, chunk in enumerate(chunks):
            embedding = self._get_embedding(chunk)
            embedding_str = json.dumps(embedding)
            cursor.execute(
                "INSERT INTO document_chunks (doc_name, chunk_index, content, embedding) VALUES (?, ?, ?, ?)",
                (doc_name, idx, chunk, embedding_str)
            )
            count += 1

        conn.commit()
        conn.close()
        logger.info("Ingested document '%s' with %d chunks into SQLite.", doc_name, count)
        return count

    # --- Cosine Similarity Search ---

    def _cosine_similarity(self, v1: List[float], v2: List[float]) -> float:
        dot_product = sum(a * b for a, b in zip(v1, v2))
        norm_a = sum(a * a for a in v1) ** 0.5
        norm_b = sum(b * b for b in v2) ** 0.5
        if not norm_a or not norm_b:
            return 0.0
        return dot_product / (norm_a * norm_b)

    def search_documents(self, query: str, doc_name: Optional[str] = None, top_k: int = 3) -> List[Dict[str, Any]]:
        """Conducts semantic similarity search inside SQLite."""
        query_vector = self._get_embedding(query)

        conn = sqlite3.connect(str(self.db_path), check_same_thread=False)
        cursor = conn.cursor()

        if doc_name:
            cursor.execute(
                "SELECT doc_name, chunk_index, content, embedding FROM document_chunks WHERE doc_name = ?",
                (doc_name,)
            )
        else:
            cursor.execute("SELECT doc_name, chunk_index, content, embedding FROM document_chunks")

        rows = cursor.fetchall()
        conn.close()

        results = []
        for name, idx, content, emb_str in rows:
            try:
                stored_vector = json.loads(emb_str)
                # Compute score if dimensions align
                if len(stored_vector) == len(query_vector):
                    score = self._cosine_similarity(query_vector, stored_vector)
                else:
                    score = 0.0
                
                # Keyword fallback weight: boost score if exact text matches
                if query.lower() in content.lower():
                    score = max(score, 0.5) + 0.15

                results.append({
                    "doc_name": name,
                    "chunk_index": idx,
                    "content": content,
                    "score": score
                })
            except Exception:
                continue

        # Sort by similarity score descending
        results.sort(key=lambda x: x["score"], reverse=True)
        return results[:top_k]
